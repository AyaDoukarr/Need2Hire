#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Script de génération du PowerPoint depuis PRESENTATION_TECHNIQUE.md
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Configuration des couleurs AXA
AXA_BLUE = RGBColor(0, 81, 186)
AXA_RED = RGBColor(226, 0, 116)
AXA_DARK = RGBColor(44, 44, 44)
GRAY_LIGHT = RGBColor(240, 240, 240)

# Lecture du fichier markdown
with open('PRESENTATION_TECHNIQUE.md', 'r', encoding='utf-8') as f:
    content = f.read()

# Parser le contenu en slides
slides_data = []
current_slide = None

for line in content.split('\n'):
    # Détecter une nouvelle slide
    slide_match = re.match(r'SLIDE (\d+) : (.+)', line)
    if slide_match:
        if current_slide:
            slides_data.append(current_slide)
        num = slide_match.group(1)
        title = slide_match.group(2).strip()
        current_slide = {'title': title, 'content': []}
    elif current_slide is not None and line.strip():
        # Ajouter le contenu (sauf les séparateurs)
        if '---' not in line:
            current_slide['content'].append(line.strip())

if current_slide:
    slides_data.append(current_slide)

# Créer le PowerPoint
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Supprimer le slide titre par défaut
if prs.slides:
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def add_title_slide(prs, title, subtitle=None):
    """Ajouter une slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre principal
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AXA_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = AXA_DARK
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    """Ajouter une slide de contenu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Bandeau titre en bleu
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AXA_BLUE
    shape.line.fill.background()
    
    # Titre de la slide
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Contenu
    content_top = Inches(1.4)
    content_height = Inches(5.8)
    
    # Créer un canvas pour le texte
    txBox = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(12.333), content_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Grouper les bullets par section
    current_section = None
    current_items = []
    
    for line in bullets:
        line = line.strip()
        if not line:
            continue
        
        # Détecter une nouvelle section (terminaison avec :)
        if line.endswith(':') and not line.startswith('-'):
            if current_section and current_items:
                # Ajouter la section précédente
                p = tf.add_paragraph()
                p.text = current_section
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = AXA_BLUE
                p.space_before = Pt(12)
                p.level = 0
                
                for item in current_items:
                    p2 = tf.add_paragraph()
                    p2.text = item
                    p2.font.size = Pt(14)
                    p2.level = 1
                    p2.space_before = Pt(4)
            
            current_section = line
            current_items = []
        elif line.startswith('-') or line.startswith('•'):
            current_items.append(line.lstrip('- ').lstrip('• '))
        else:
            # Ligne de transition ou numérotée
            if line[0].isdigit() and '. ' in line:
                # C'est un élément numéroté
                if current_section and current_items:
                    p = tf.add_paragraph()
                    p.text = current_section
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = AXA_BLUE
                    p.space_before = Pt(12)
                    
                    for item in current_items:
                        p2 = tf.add_paragraph()
                        p2.text = item
                        p2.font.size = Pt(14)
                        p2.level = 1
                        p2.space_before = Pt(4)
                    
                    current_items = []
                    current_section = None
                
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = AXA_DARK
                p.space_before = Pt(10)
            else:
                current_items.append(line)
    
    # Ajouter la dernière section
    if current_section or current_items:
        if current_section:
            p = tf.add_paragraph()
            p.text = current_section
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = AXA_BLUE
            p.space_before = Pt(12)
            
            for item in current_items:
                p2 = tf.add_paragraph()
                p2.text = item
                p2.font.size = Pt(14)
                p2.level = 1
                p2.space_before = Pt(4)
        else:
            for item in current_items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(14)
                p.level = 0
                p.space_before = Pt(4)
    
    return slide

# Slide de titre
add_title_slide(prs, "IA DE QUALIFICATION DU BESOIN DE RECRUTEMENT", 
               "SQORUS | IA RH")

# Slides de contenu
for slide_data in slides_data:
    title = slide_data['title']
    content = slide_data['content']
    add_content_slide(prs, title, content)

# Enregistrer le fichier
output_path = 'need2hire_presentation.pptx'
prs.save(output_path)
print(f"Fichier PowerPoint créé : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")